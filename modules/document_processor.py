"""
document_processor.py
----------------------
Handles extraction of raw text from the file types the app supports:
PDF (PyMuPDF), DOCX (python-docx), PPTX (python-pptx), and plain TXT.

Every extractor returns a plain string. Errors are caught and re-raised
as a single DocumentProcessingError so the UI layer only has to handle
one exception type.
"""

from __future__ import annotations
import io
import fitz  # PyMuPDF
import docx
from pptx import Presentation


class DocumentProcessingError(Exception):
    """Raised when a file cannot be read or parsed."""
    pass


def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract text from a PDF file's raw bytes using PyMuPDF."""
    try:
        text_parts = []
        with fitz.open(stream=file_bytes, filetype="pdf") as pdf:
            for page_num, page in enumerate(pdf, start=1):
                page_text = page.get_text("text")
                if page_text.strip():
                    text_parts.append(f"\n--- Page {page_num} ---\n{page_text}")
        text = "".join(text_parts).strip()
        if not text:
            raise DocumentProcessingError(
                "No extractable text found in this PDF (it may be scanned/image-only)."
            )
        return text
    except DocumentProcessingError:
        raise
    except Exception as exc:
        raise DocumentProcessingError(f"Failed to read PDF: {exc}") from exc


def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract text (paragraphs + tables) from a DOCX file's raw bytes."""
    try:
        document = docx.Document(io.BytesIO(file_bytes))
        text_parts = [p.text for p in document.paragraphs if p.text.strip()]

        # Include table content too, since study notes often use tables.
        for table in document.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)

        text = "\n".join(text_parts).strip()
        if not text:
            raise DocumentProcessingError("No extractable text found in this DOCX file.")
        return text
    except DocumentProcessingError:
        raise
    except Exception as exc:
        raise DocumentProcessingError(f"Failed to read DOCX: {exc}") from exc


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract text from every slide (shapes + notes) of a PPTX file."""
    try:
        presentation = Presentation(io.BytesIO(file_bytes))
        text_parts = []
        for slide_num, slide in enumerate(presentation.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    slide_text.append(shape.text_frame.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            slide_text.append(row_text)
            # Speaker notes often contain valuable explanations.
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                slide_text.append(f"(Notes: {slide.notes_slide.notes_text_frame.text})")

            if slide_text:
                text_parts.append(f"\n--- Slide {slide_num} ---\n" + "\n".join(slide_text))

        text = "".join(text_parts).strip()
        if not text:
            raise DocumentProcessingError("No extractable text found in this PPTX file.")
        return text
    except DocumentProcessingError:
        raise
    except Exception as exc:
        raise DocumentProcessingError(f"Failed to read PPTX: {exc}") from exc


def extract_text_from_txt(file_bytes: bytes) -> str:
    """Decode a plain text file, trying a couple of common encodings."""
    for encoding in ("utf-8", "latin-1"):
        try:
            text = file_bytes.decode(encoding).strip()
            if not text:
                raise DocumentProcessingError("The uploaded TXT file is empty.")
            return text
        except UnicodeDecodeError:
            continue
    raise DocumentProcessingError("Could not decode this TXT file's encoding.")


def extract_text(filename: str, file_bytes: bytes) -> str:
    """
    Dispatch to the correct extractor based on file extension.
    Raises DocumentProcessingError for unsupported types or failures.
    """
    extension = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""

    extractors = {
        "pdf": extract_text_from_pdf,
        "docx": extract_text_from_docx,
        "pptx": extract_text_from_pptx,
        "txt": extract_text_from_txt,
    }

    if extension not in extractors:
        raise DocumentProcessingError(
            f"Unsupported file type '.{extension}'. Please upload PDF, DOCX, PPTX, or TXT."
        )

    return extractors[extension](file_bytes)
